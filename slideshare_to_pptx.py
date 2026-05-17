#!/usr/bin/env python3
from __future__ import annotations

import argparse
import html
import io
import json
import re
import tempfile
import warnings
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches
from requests.adapters import HTTPAdapter
from urllib3.exceptions import NotOpenSSLWarning
from urllib3.util.retry import Retry

warnings.filterwarnings("ignore", category=NotOpenSSLWarning)

DEFAULT_OUTPUT_DIR = "output"
DEFAULT_TIMEOUT = 30
HIGH_RES_BUCKET = 75
LOW_RES_BUCKET = 85
DOWNLOAD_SIZES = (2048, 1600, 1200, 1024, 638, 320)
USER_AGENT = (
    "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
    "AppleWebKit/537.36 (KHTML, like Gecko) "
    "Chrome/136.0.0.0 Safari/537.36"
)

IMAGE_URL_RE = re.compile(r"https://image\.slidesharecdn\.com/[^\"'<> ]+")
SLIDE_URL_RE = re.compile(
    r"(?P<prefix>https://image\.slidesharecdn\.com/.+?)/(?P<bucket>\d+)/"
    r"(?P<basename>.+?)-(?P<slide>\d+)-(?P<size>\d+)\.(?P<ext>jpg|jpeg|png|webp)"
    r"(?P<query>\?[^\"'<> ]*)?$",
    re.IGNORECASE,
)
TITLE_SUFFIX_RE = re.compile(
    r"\s*(?:\|\s*(?:PPTX?|PDF)|[-|]\s*SlideShare)\s*$",
    re.IGNORECASE,
)


class SlideShareError(RuntimeError):
    pass


@dataclass(frozen=True)
class SlideCandidate:
    slide_number: int
    size: int
    url: str


@dataclass(frozen=True)
class ExportResult:
    pptx_path: Path
    slide_count: int
    title: str


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Download public SlideShare slide images and package them into a PPTX."
    )
    parser.add_argument("url", help="SlideShare slideshow URL")
    parser.add_argument(
        "--output-dir",
        default=DEFAULT_OUTPUT_DIR,
        help="Directory where images and the generated PPTX will be written",
    )
    parser.add_argument(
        "--title",
        default=None,
        help="Optional override for the generated file name",
    )
    parser.add_argument(
        "--timeout",
        type=int,
        default=DEFAULT_TIMEOUT,
        help="HTTP timeout in seconds",
    )
    parser.add_argument(
        "--keep-webp",
        action="store_true",
        help="Keep raw SlideShare image responses alongside the converted PNG files",
    )
    return parser.parse_args()


def build_session() -> requests.Session:
    retry = Retry(
        total=3,
        backoff_factor=0.6,
        status_forcelist=(429, 500, 502, 503, 504),
        allowed_methods=frozenset({"GET", "HEAD"}),
        raise_on_status=False,
    )
    adapter = HTTPAdapter(max_retries=retry)
    session = requests.Session()
    session.headers.update(
        {
            "User-Agent": USER_AGENT,
            "Accept": (
                "text/html,application/xhtml+xml,application/xml;q=0.9,"
                "image/avif,image/webp,*/*;q=0.8"
            ),
            "Accept-Language": "en-US,en;q=0.9",
        }
    )
    session.mount("http://", adapter)
    session.mount("https://", adapter)
    return session


def normalize_slideshare_url(url: str) -> str:
    value = url.strip()
    if not value:
        raise SlideShareError("SlideShare URL is required")
    if "://" not in value:
        value = f"https://{value}"
    parsed = urlparse(value)
    if parsed.scheme not in {"http", "https"}:
        raise SlideShareError("URL must use http or https")
    hostname = (parsed.netloc or "").lower()
    if "slideshare.net" not in hostname:
        raise SlideShareError("Only public SlideShare URLs are supported")
    return value


def sanitize_filename(value: str) -> str:
    cleaned = html.unescape(value).strip()
    cleaned = re.sub(r"[\\/:*?\"<>|]+", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)
    cleaned = cleaned.rstrip(". ")
    return cleaned or "slideshare-export"


def clean_title(value: str) -> str:
    title = html.unescape(value).strip()
    while True:
        updated = TITLE_SUFFIX_RE.sub("", title).strip()
        if updated == title:
            return updated
        title = updated


def fetch_page(session: requests.Session, url: str, timeout: int) -> str:
    normalized_url = normalize_slideshare_url(url)
    response = session.get(normalized_url, timeout=timeout)
    response.raise_for_status()
    resolved_host = urlparse(response.url).netloc.lower()
    if "slideshare.net" not in resolved_host:
        raise SlideShareError(f"unexpected redirect target: {response.url}")
    if "text/html" not in response.headers.get("content-type", ""):
        raise SlideShareError("expected an HTML SlideShare page")
    return response.text


def iter_json_ld_objects(soup: BeautifulSoup) -> Iterable[dict]:
    for script in soup.select('script[type="application/ld+json"]'):
        text = script.string or script.get_text(strip=True)
        if not text:
            continue
        try:
            payload = json.loads(text)
        except json.JSONDecodeError:
            continue
        if isinstance(payload, dict):
            yield payload
        elif isinstance(payload, list):
            for item in payload:
                if isinstance(item, dict):
                    yield item


def extract_title(html_text: str) -> str:
    soup = BeautifulSoup(html_text, "html.parser")
    selectors = (
        ('meta[name="title"]', "content"),
        ('meta[property="og:title"]', "content"),
        ('meta[name="twitter:title"]', "content"),
        ("title", None),
    )

    for selector, attr in selectors:
        element = soup.select_one(selector)
        if not element:
            continue
        value = element.get(attr) if attr else element.get_text(strip=True)
        if value:
            cleaned = clean_title(value)
            if cleaned:
                return cleaned

    for obj in iter_json_ld_objects(soup):
        name = obj.get("name")
        if isinstance(name, str):
            cleaned = clean_title(name)
            if cleaned:
                return cleaned

    return "slideshare-export"


def parse_slide_candidate(url: str) -> SlideCandidate | None:
    match = SLIDE_URL_RE.match(url)
    if not match:
        return None
    return SlideCandidate(
        slide_number=int(match.group("slide")),
        size=int(match.group("size")),
        url=url,
    )


def validate_slide_sequence(slide_numbers: list[int]) -> None:
    if not slide_numbers:
        raise SlideShareError("no slide image URLs were found in the page HTML")

    expected = list(range(1, slide_numbers[-1] + 1))
    if slide_numbers != expected:
        missing = sorted(set(expected) - set(slide_numbers))
        raise SlideShareError(
            "slide sequence is incomplete; missing slide numbers: "
            + ", ".join(str(value) for value in missing)
        )


def extract_slide_urls(html_text: str) -> dict[int, str]:
    best_by_slide: dict[int, SlideCandidate] = {}

    for url in sorted(set(IMAGE_URL_RE.findall(html_text))):
        candidate = parse_slide_candidate(url)
        if candidate is None:
            continue
        existing = best_by_slide.get(candidate.slide_number)
        if existing is None or candidate.size > existing.size:
            best_by_slide[candidate.slide_number] = candidate

    slide_numbers = sorted(best_by_slide)
    validate_slide_sequence(slide_numbers)
    return {
        slide_number: best_by_slide[slide_number].url
        for slide_number in slide_numbers
    }


def build_download_candidates(url: str) -> list[str]:
    match = SLIDE_URL_RE.match(url)
    if not match:
        return [url]

    prefix = match.group("prefix")
    basename = match.group("basename")
    slide = match.group("slide")
    ext = match.group("ext")
    current_size = int(match.group("size"))

    candidate_sizes = list(dict.fromkeys((*DOWNLOAD_SIZES, current_size)))
    ordered_sizes = sorted(candidate_sizes, reverse=True)

    urls: list[str] = []
    for size in ordered_sizes:
        bucket = HIGH_RES_BUCKET if size > 320 else LOW_RES_BUCKET
        urls.append(f"{prefix}/{bucket}/{basename}-{slide}-{size}.{ext}")
    urls.append(url)
    return list(dict.fromkeys(urls))


def validate_image_bytes(image_bytes: bytes) -> None:
    with Image.open(io.BytesIO(image_bytes)) as image:
        image.load()


def fetch_slide_image(
    session: requests.Session,
    preview_url: str,
    timeout: int,
) -> tuple[bytes, str, str]:
    last_error: Exception | None = None

    for candidate_url in build_download_candidates(preview_url):
        try:
            response = session.get(candidate_url, timeout=timeout)
            if not response.ok:
                continue
            content_type = response.headers.get("content-type", "")
            if not content_type.startswith("image/"):
                continue
            validate_image_bytes(response.content)
            return response.content, candidate_url, content_type
        except (requests.RequestException, OSError, UnidentifiedImageError) as exc:
            last_error = exc

    if last_error is not None:
        raise SlideShareError(f"failed to download a valid slide image: {last_error}") from last_error
    raise SlideShareError("failed to download a valid slide image")


def convert_to_png(image_bytes: bytes, output_path: Path) -> None:
    with Image.open(io.BytesIO(image_bytes)) as image:
        if image.mode not in ("RGB", "RGBA"):
            image = image.convert("RGBA" if "A" in image.getbands() else "RGB")
        image.save(output_path, format="PNG")


def guess_extension(content_type: str) -> str:
    mapping = {
        "image/webp": ".webp",
        "image/jpeg": ".jpg",
        "image/png": ".png",
    }
    return mapping.get(content_type.lower(), ".bin")


def clean_previous_artifacts(output_root: Path, pptx_path: Path) -> None:
    images_dir = output_root / "images"
    if images_dir.exists():
        for path in images_dir.glob("slide-*.*"):
            path.unlink()
    pptx_path.unlink(missing_ok=True)


def download_slides(
    session: requests.Session,
    slide_urls: dict[int, str],
    images_dir: Path,
    timeout: int,
    keep_webp: bool,
) -> list[Path]:
    images_dir.mkdir(parents=True, exist_ok=True)
    png_paths: list[Path] = []

    for slide_number, preview_url in slide_urls.items():
        image_bytes, resolved_url, content_type = fetch_slide_image(
            session=session,
            preview_url=preview_url,
            timeout=timeout,
        )

        if keep_webp:
            raw_extension = guess_extension(content_type)
            raw_path = images_dir / f"slide-{slide_number:02d}{raw_extension}"
            raw_path.write_bytes(image_bytes)

        png_path = images_dir / f"slide-{slide_number:02d}.png"
        convert_to_png(image_bytes, png_path)
        png_paths.append(png_path)
        print(f"downloaded slide {slide_number:02d} from {resolved_url}")

    return png_paths


def get_slide_dimensions(image_path: Path) -> tuple[int, int]:
    with Image.open(image_path) as image:
        return image.size


def build_pptx(image_paths: Iterable[Path], output_path: Path) -> None:
    ordered_paths = list(image_paths)
    if not ordered_paths:
        raise SlideShareError("no image files were downloaded")

    width_px, height_px = get_slide_dimensions(ordered_paths[0])
    aspect_ratio = height_px / width_px
    slide_width = Inches(13.333)
    slide_height = int(slide_width * aspect_ratio)

    presentation = Presentation()
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height
    blank_layout = presentation.slide_layouts[6]

    for image_path in ordered_paths:
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

    presentation.save(output_path)


def verify_generated_pptx(output_path: Path, expected_slide_count: int) -> None:
    if not output_path.exists() or output_path.stat().st_size == 0:
        raise SlideShareError("generated PPTX is missing or empty")

    with zipfile.ZipFile(output_path) as archive:
        names = archive.namelist()
        slide_entries = sorted(
            name
            for name in names
            if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        )
        media_entries = sorted(
            name
            for name in names
            if re.fullmatch(r"ppt/media/image\d+\.(png|jpg|jpeg|webp)", name)
        )

    if len(slide_entries) != expected_slide_count:
        raise SlideShareError(
            f"generated PPTX has {len(slide_entries)} slides; expected {expected_slide_count}"
        )
    if len(media_entries) < expected_slide_count:
        raise SlideShareError(
            f"generated PPTX only contains {len(media_entries)} media files for "
            f"{expected_slide_count} slides"
        )


def write_source_note(output_root: Path, source_url: str) -> None:
    note_path = output_root / "SOURCE.txt"
    note_path.write_text(
        "Source URL:\n"
        f"{source_url}\n\n"
        "This PPTX is assembled from the publicly exposed slide images returned by "
        "SlideShare.\n"
        "It does not recover the original editable deck structure.\n",
        encoding="utf-8",
    )


def export_slideshare_pptx(
    url: str,
    output_dir: str | Path = DEFAULT_OUTPUT_DIR,
    title: str | None = None,
    timeout: int = DEFAULT_TIMEOUT,
    keep_webp: bool = False,
    session: requests.Session | None = None,
) -> ExportResult:
    active_session = session or build_session()
    html_text = fetch_page(active_session, url, timeout)
    safe_title = sanitize_filename(title or extract_title(html_text))
    slide_urls = extract_slide_urls(html_text)

    output_root = Path(output_dir).resolve() / safe_title
    images_dir = output_root / "images"
    pptx_path = output_root / f"{safe_title}.pptx"

    output_root.mkdir(parents=True, exist_ok=True)
    clean_previous_artifacts(output_root, pptx_path)

    image_paths = download_slides(
        session=active_session,
        slide_urls=slide_urls,
        images_dir=images_dir,
        timeout=timeout,
        keep_webp=keep_webp,
    )
    build_pptx(image_paths, pptx_path)
    verify_generated_pptx(pptx_path, len(image_paths))
    write_source_note(output_root, normalize_slideshare_url(url))
    return ExportResult(
        pptx_path=pptx_path,
        slide_count=len(image_paths),
        title=safe_title,
    )


def export_slideshare_pptx_to_tempfile(
    url: str,
    title: str | None = None,
    timeout: int = DEFAULT_TIMEOUT,
    session: requests.Session | None = None,
) -> ExportResult:
    temp_root = Path(tempfile.mkdtemp(prefix="slideshare-export-"))
    return export_slideshare_pptx(
        url=url,
        output_dir=temp_root,
        title=title,
        timeout=timeout,
        keep_webp=False,
        session=session,
    )


def main() -> int:
    args = parse_args()
    try:
        result = export_slideshare_pptx(
            url=args.url,
            output_dir=args.output_dir,
            title=args.title,
            timeout=args.timeout,
            keep_webp=args.keep_webp,
        )
    except (requests.RequestException, OSError, SlideShareError) as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    print(f"saved pptx: {result.pptx_path}")
    print(f"slides found: {result.slide_count}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
